"""
File extractor for various document formats
"""

import asyncio
from pathlib import Path
from typing import Optional
import re

class FileExtractor:
    """Extract text from various file formats"""

    def __init__(self):
        self.supported_formats = {
            '.txt': self._extract_text_file,
            '.md': self._extract_markdown_file,
            '.pdf': self._extract_pdf_file,
            '.docx': self._extract_docx_file,
            '.doc': self._extract_doc_file,
            '.pptx': self._extract_pptx_file,
            '.ppt': self._extract_ppt_file,
            '.xlsx': self._extract_xlsx_file,
            '.xls': self._extract_xls_file,
            '.py': self._extract_code_file,
            '.js': self._extract_code_file,
            '.java': self._extract_code_file,
            '.cpp': self._extract_code_file,
            '.c': self._extract_code_file,
            '.html': self._extract_text_file,
            '.css': self._extract_text_file,
            '.json': self._extract_text_file,
            '.xml': self._extract_text_file,
        }

    async def extract_text(self, file_path: str) -> str:
        """Extract text from a file based on its extension"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = path.suffix.lower()

        if extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {extension}")

        try:
            extractor = self.supported_formats[extension]
            # Run extraction in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            text = await loop.run_in_executor(None, extractor, file_path)
            return self._clean_text(text)
        except Exception as e:
            raise Exception(f"Failed to extract text from {file_path}: {e}")

    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _extract_markdown_file(self, file_path: str) -> str:
        """Extract text from markdown files, preserving structure"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        # Remove markdown formatting but preserve structure
        # Remove headers formatting but keep text
        content = re.sub(r'^#{1,6}\s+', '', content, flags=re.MULTILINE)
        # Remove emphasis markers
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
        content = re.sub(r'\*([^*]+)\*', r'\1', content)
        content = re.sub(r'_([^_]+)_', r'\1', content)
        # Remove links but keep text
        content = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', content)
        # Remove code blocks markers but keep content
        content = re.sub(r'```[^\n]*\n', '', content)
        content = re.sub(r'```', '', content)
        # Remove inline code markers
        content = re.sub(r'`([^`]+)`', r'\1', content)

        return content

    def _extract_pdf_file(self, file_path: str) -> str:
        """Extract text from PDF files"""
        try:
            import PyPDF2
        except ImportError:
            raise ImportError("PyPDF2 not installed. Install with: pip install PyPDF2")

        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"

        return text

    def _extract_docx_file(self, file_path: str) -> str:
        """Extract text from DOCX files"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")

        doc = Document(file_path)
        text = ""

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"

        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"

        return text

    def _extract_doc_file(self, file_path: str) -> str:
        """Extract text from DOC files (limited support)"""
        # DOC files are complex binary format
        # For now, just return a message that conversion is needed
        return "DOC file format requires external conversion to DOCX for text extraction. Please convert to DOCX format."

    def _extract_pptx_file(self, file_path: str) -> str:
        """Extract text from PPTX files"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        prs = Presentation(file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"

        return text

    def _extract_ppt_file(self, file_path: str) -> str:
        """Extract text from PPT files (limited support)"""
        return "PPT file format requires external conversion to PPTX for text extraction. Please convert to PPTX format."

    def _extract_xlsx_file(self, file_path: str) -> str:
        """Extract text from XLSX files"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl not installed. Install with: pip install openpyxl")

        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ""

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"

            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text += row_text + "\n"

            text += "\n"

        return text

    def _extract_xls_file(self, file_path: str) -> str:
        """Extract text from XLS files (limited support)"""
        try:
            import xlrd
        except ImportError:
            raise ImportError("xlrd not installed. Install with: pip install xlrd")

        workbook = xlrd.open_workbook(file_path)
        text = ""

        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            text += f"Sheet: {sheet.name}\n"

            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                row_text = " | ".join(str(cell) for cell in row_values if str(cell).strip())
                if row_text.strip():
                    text += row_text + "\n"

            text += "\n"

        return text

    def _extract_code_file(self, file_path: str) -> str:
        """Extract text from code files, preserving structure"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        # For code files, we want to preserve comments and structure
        # but clean up excessive whitespace
        lines = content.split('\n')
        cleaned_lines = []

        for line in lines:
            # Skip completely empty lines if there are too many consecutive
            if line.strip() or (cleaned_lines and cleaned_lines[-1].strip()):
                cleaned_lines.append(line)

        return '\n'.join(cleaned_lines)

    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        if not text:
            return ""

        # Normalize whitespace
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)  # Multiple newlines to double
        text = re.sub(r'[ \t]+', ' ', text)  # Multiple spaces to single

        # Remove excessive whitespace at start/end
        text = text.strip()

        # Ensure we have some content
        if len(text) < 10:
            return "Text extracted but content appears to be minimal or corrupted."

        return text

    def get_supported_formats(self) -> list:
        """Get list of supported file extensions"""
        return list(self.supported_formats.keys())

    def is_format_supported(self, file_path: str) -> bool:
        """Check if a file format is supported"""
        return Path(file_path).suffix.lower() in self.supported_formats

